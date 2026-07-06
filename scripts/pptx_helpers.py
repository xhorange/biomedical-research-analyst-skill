"""Reusable python-pptx building blocks for biomedical report decks, so each
new deck doesn't reinvent shape/textframe boilerplate. Import what you need;
these are plain functions operating on a `pptx.Presentation` / `Slide`, not a
framework you have to adopt wholesale.

Typical setup:

    from pptx import Presentation
    from pptx.util import Inches
    from pptx_helpers import NAVY, ACCENT, add_slide, add_title_bar, add_pic_fit, add_bullets, add_table

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    s = add_slide(prs, BLANK)
    add_title_bar(s, prs, "Slide title", "optional subtitle", page_no=1, total=10)
    ...
    prs.save("output.pptx")

After building, ALWAYS read the .pptx back with python-pptx and check slide
count / picture / table presence before calling the deck done - see
references/ppt_structure.md.
"""
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)

DEFAULT_FONT = "Hiragino Sans GB"  # swap for a font that covers your target language/script


def add_slide(prs, blank_layout):
    return prs.slides.add_slide(blank_layout)


def add_bg(slide, sw, sh, color=RGBColor(0xFF, 0xFF, 0xFF)):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)
    return rect


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT, font=DEFAULT_FONT,
                 anchor=None, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return tb


def add_title_bar(slide, prs, title, subtitle=None, page_no=None, total=None,
                   bar_color=NAVY, bar_height_in=1.05):
    from pptx.util import Inches
    sw, sh = prs.slide_width, prs.slide_height
    add_bg(slide, sw, sh)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(bar_height_in))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_textbox(slide, Inches(0.5), Inches(0.12), sw - Inches(1), Inches(0.6), title,
                size=26, bold=True, color=RGBColor(255, 255, 255))
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.62), sw - Inches(1), Inches(0.4), subtitle,
                    size=13, color=RGBColor(210, 220, 235))
    if page_no is not None:
        add_textbox(slide, sw - Inches(1.1), sh - Inches(0.45), Inches(0.8), Inches(0.3),
                    f"{page_no}/{total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)


def add_pic_fit(slide, path, left, top, max_w, max_h):
    """Place an image scaled to fit inside a box, preserving aspect ratio and
    centering it - avoids the stretched/distorted look of a naive width/height set."""
    from PIL import Image
    with Image.open(path) as im:
        w_px, h_px = im.size
    ratio = w_px / h_px
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        w = max_w
        h = int(max_w / ratio)
    else:
        h = max_h
        w = int(max_h * ratio)
    l = left + int((max_w - w) / 2)
    t = top + int((max_h - h) / 2)
    slide.shapes.add_picture(path, l, t, width=w, height=h)


def add_bullets(slide, left, top, width, height, items, size=15,
                 color=RGBColor(0x20, 0x20, 0x20), font=DEFAULT_FONT, space_after=8):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.space_after = Pt(space_after)
        p.line_spacing = 1.2
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = font
    return tb


def add_table(slide, left, top, width, height, header, rows, col_widths=None,
              header_color=NAVY, font_size=11):
    n_rows = len(rows) + 1
    n_cols = len(header)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, htext in enumerate(header):
        cell = table.cell(0, j)
        cell.text = str(htext)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(font_size)
                r.font.color.rgb = RGBColor(255, 255, 255)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 == 0 else RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
    return table


def verify_deck(pptx_path, expected_slides=None):
    """Read a built deck back and print a structural summary - call this
    before telling the user the deck is done. Returns the list of
    (pics, tables) counts per slide for programmatic checks if needed."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    n = len(prs.slides)
    print(f"{pptx_path}: {n} slides")
    if expected_slides is not None and n != expected_slides:
        print(f"WARNING: expected {expected_slides} slides, got {n}")
    counts = []
    for i, slide in enumerate(prs.slides, 1):
        n_pics = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        n_tables = sum(1 for sh in slide.shapes if sh.has_table)
        counts.append((n_pics, n_tables))
        print(f"  slide {i}: pics={n_pics} tables={n_tables}")
    return counts
